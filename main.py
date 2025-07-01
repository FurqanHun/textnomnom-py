import os
import sys
import re
import logging
import subprocess
from PyPDF2 import PdfReader
from pptx import Presentation
import shutil
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
import io
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.service import Service as ChromeService
from selenium.webdriver.firefox.service import Service as FirefoxService
from selenium.webdriver.chrome.options import Options as ChromeOptions
from selenium.webdriver.firefox.options import Options as FirefoxOptions
from markdownify import markdownify as md

# if you have installed the driver, update this path in and make sure to use doble quotes(""), else leave it as None
GECKO_DRIVER_PATH = "/mnt/qanhdd/some-stuff/geckodriver"
CHROME_DRIVER_PATH = None

CHROMIUM_BASED_BROWSER_PATH = None
FIREFOX_BASED_BROWSER_PATH = None

# Configure logging
LOG_FILE = "file_processing.log"
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler(LOG_FILE, encoding="utf-8"),
        logging.StreamHandler()  # Also log to console
    ]
)
def clear_log():
    """Clear the content of the log file."""
    try:
        with open(LOG_FILE, "w", encoding="utf-8") as file:
            file.truncate(0)  # Truncate the file to 0 bytes
        logging.info(f"Log file '{LOG_FILE}' cleared.")
    except Exception as e:
        logging.error(f"Error clearing log file: {e}")

#URL part
def get_driver():
    # Check for Chrome availability
    if CHROMIUM_BASED_BROWSER_PATH != None:
        chrome_path = CHROMIUM_BASED_BROWSER_PATH
    else:
        if sys.platform == "win32":
            chrome_path = "C:\\Program Files\\Google\\Chrome\\Application\\chrome.exe"
        else:
            chrome_path = "/usr/bin/google-chrome"

    if os.path.exists(chrome_path) and CHROME_DRIVER_PATH:
        chrome_options = ChromeOptions()
        chrome_options.binary_location = chrome_path
        chrome_options.add_argument("--headless")
        chrome_service = ChromeService(executable_path=CHROME_DRIVER_PATH)
        driver = webdriver.Chrome(service=chrome_service, options=chrome_options)
        return driver, "chrome"

    # Check for Firefox availability
    if FIREFOX_BASED_BROWSER_PATH != None:
        firefox_path = FIREFOX_BASED_BROWSER_PATH
    else:
        if sys.platform == "win32":
            firefox_path = "C:\\Program Files\\Mozilla Firefox\\firefox.exe"
        else:
            firefox_path = "/usr/bin/mullvad-browser"

    if os.path.exists(firefox_path) and GECKO_DRIVER_PATH:
        firefox_options = FirefoxOptions()
        firefox_options.binary_location = firefox_path
        firefox_options.add_argument("--headless")
        firefox_service = FirefoxService(executable_path=GECKO_DRIVER_PATH)
        driver = webdriver.Firefox(service=firefox_service, options=firefox_options)
        return driver, "firefox"

    raise EnvironmentError("Neither Chrome nor Firefox is installed on this system.")

def sanitize_filename(filename):
    return re.sub(r'[<>:"/\\|?*]', '_', filename)

def scrape_and_save(url):
    logging.info(f"Scraping URL: {url}")
    try:
        driver, browser = get_driver()
        driver.get(url)
        if "pdf" in driver.current_url:
            logging.info("PDF file detected. Saving PDF...")
            pdf_path = os.path.join(os.getcwd(), "temp.pdf")
            driver.get(driver.current_url)
            with open(pdf_path, "wb") as file:
                file.write(driver.page_source.encode("utf-8"))
            driver.quit()
            logging.info(f"PDF saved to {pdf_path}")
            return
        logging.info(f"Page title: {driver.title}")

    except Exception as e:
        logging.error(f"Error scraping URL: {e}")
        return

    body_element = driver.find_element(By.TAG_NAME, "body")
    html_content = body_element.get_attribute("outerHTML")

    logging.info("Converting HTML to Markdown...")
    markdown_content = md(html_content)
    page_title = driver.title.strip().replace(" ", "_").replace("/", "_")
    sanitized_title = sanitize_filename(page_title)
    filename = f"{sanitized_title}.md"

    logging.info(f"Saving Markdown content to {filename}...")
    with open(filename, "w", encoding="utf-8") as file:
        file.write(markdown_content)
    driver.quit()
    logging.info(f"Markdown content saved to {filename} using {browser.capitalize()}.")

def extract_text_from_image(file_path):
    """Extract text from image using Tesseract OCR."""
    try:
        text = pytesseract.image_to_string(Image.open(file_path))
        logging.info(f"Extracted text from image: {file_path}")
        return text.strip()
    except Exception as e:
        logging.error(f"Error processing image file: {file_path}, Error: {e}")
        return None

def extract_text_from_pdf(file_path, trigger_ocr, ocr_mix=False):
    """Extract text from each page of the PDF, optionally with OCR for images."""
    pdf_text = []
    try:
        with open(file_path, "rb") as file:
            reader = PdfReader(file)

            for page_num, page in enumerate(reader.pages):
                # Regular text extraction
                text = page.extract_text()
                if text:
                    pdf_text.append(text)

                # OCR extraction if triggered or ocr_mix is True
                if trigger_ocr or ocr_mix:
                    try:
                        images = convert_from_path(file_path, first_page=page_num + 1, last_page=page_num + 1)
                        if images:
                            image = images[0]
                            ocr_text = pytesseract.image_to_string(image)
                            if ocr_text and ocr_text.strip() not in pdf_text: # Avoid duplicate text
                                pdf_text.append(f"[OCR from page {page_num + 1}]\n{ocr_text.strip()}")
                    except Exception as e:
                        logging.warning(f"Error during OCR on page {page_num + 1}: {e}")

        logging.info(f"Extracted text from PDF: {file_path} (OCR Mix: {ocr_mix})")
    except Exception as e:
        logging.error(f"Error reading the PDF file: {e}")
        return None
    return "\n".join(pdf_text)


def convert_ppt_to_pptx(ppt_path):
    """Convert .ppt to .pptx using external tools."""
    temp_dir = os.path.dirname(ppt_path)
    temp_filename = f"temp_{os.path.basename(ppt_path)}x"
    pptx_path = os.path.join(temp_dir, temp_filename)

    if sys.platform == 'win32':
        try:
            import win32com.client
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            # powerpoint.Visible = False
            deck = powerpoint.Presentations.Open(ppt_path)
            deck.SaveAs(pptx_path, 24)  # 24 = .pptx format
            deck.Close()
            powerpoint.Quit()
            logging.info(f"Successfully converted .ppt to .pptx using PowerPoint: {pptx_path}")
            return pptx_path
        except Exception as e:
            logging.error(f"PowerPoint automation failed: {e}")
    else:
        converters = [
            {
                "name": "soffice",
                "cmd": ["soffice", "--headless", "--convert-to", "pptx", "--outdir", temp_dir, ppt_path],
            },
            {
                "name": "unoconv",
                "cmd": ["unoconv", "-f", "pptx", "-o", temp_dir, ppt_path],
            }
        ]

        for converter in converters:
            if shutil.which(converter["cmd"][0]):
                try:
                    logging.info(f"Attempting .ppt to .pptx conversion using {converter['name']}...")
                    result = subprocess.run(converter["cmd"], capture_output=True, text=True)
                    if result.returncode == 0 and os.path.exists(pptx_path):
                        logging.info(f"Successfully converted .ppt to .pptx: {pptx_path}")
                        return pptx_path
                    else:
                        logging.warning(f"{converter['name']} failed. Output: {result.stderr.strip()}")
                except Exception as e:
                    logging.error(f"{converter['name']} conversion failed: {e}")
        logging.error("No conversion tools succeeded. Please convert the .ppt manually.")
        return None

def extract_text_from_pptx(file_path, trigger_ocr, ocr_mix=False):
    """Extract text from PowerPoint presentation, applying OCR if necessary."""
    pptx_text = []
    try:
        prs = Presentation(file_path)
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []

            for shape in slide.shapes:
                # Check if the shape contains text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # If OCR is triggered or ocr_mix is True and the shape contains an image, process it
                elif (trigger_ocr or ocr_mix) and hasattr(shape, "image"):
                    logging.info(f"Found image in Slide {slide_number}, triggering OCR...")
                    image = shape.image
                    image_bytes = image.blob
                    image_pil = Image.open(io.BytesIO(image_bytes))
                    text = pytesseract.image_to_string(image_pil)
                    if text:
                        slide_text.append(f"[OCR from Slide {slide_number}]\n{text.strip()}")

            if slide_text:
                pptx_text.append(f"[Slide {slide_number}]\n" + "\n".join(slide_text))

        logging.info(f"Extracted text from PowerPoint: {file_path} (OCR Mix: {ocr_mix})")
    except Exception as e:
        logging.error(f"Error reading the PowerPoint file: {e}")
        return None
    return "\n\n".join(pptx_text)

def convert_pptx_to_pdf(pptx_path):
    """Convert .pptx to .pdf using external tools."""
    temp_dir = os.path.dirname(pptx_path)
    pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"

    if sys.platform == 'win32':
        try:
            import win32com.client
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            # powerpoint.Visible = False
            presentation = powerpoint.Presentations.Open(pptx_path)
            presentation.SaveAs(pdf_path, FileFormat=32)  # 32 = ppSaveAsPDF
            presentation.Close()
            powerpoint.Quit()
            logging.info(f"Successfully converted .pptx to .pdf: {pdf_path}")
            return pdf_path
        except Exception as e:
            logging.error(f"PowerPoint automation failed: {e}")
            return None
    else:
        converters = [
            {
                "name": "soffice",
                "cmd": ["soffice", "--headless", "--convert-to", "pdf", "--outdir", temp_dir, pptx_path],
            },
            {
                "name": "unoconv",
                "cmd": ["unoconv", "-f", "pdf", "-o", temp_dir, pptx_path],
            }
        ]

        for converter in converters:
            if shutil.which(converter["cmd"][0]):
                try:
                    logging.info(f"Attempting .pptx to .pdf conversion using {converter['name']}...")
                    result = subprocess.run(converter["cmd"], capture_output=True, text=True)
                    if result.returncode == 0 and os.path.exists(pdf_path):
                        logging.info(f"Successfully converted .pptx to .pdf: {pdf_path}")
                        return pdf_path
                    else:
                        logging.warning(f"{converter['name']} failed. Output: {result.stderr.strip()}")
                except Exception as e:
                    logging.error(f"{converter['name']} conversion failed: {e}")
        logging.error("No conversion tools succeeded. Please convert the .pptx manually.")
        return None

def extract_text_from_file(file_path, trigger_ocr, ocr_mix=False):
    """Extract text based on file type."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path, trigger_ocr, ocr_mix)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path, trigger_ocr, ocr_mix)
    elif ext == ".ppt":
        logging.info(".ppt file detected. Attempting conversion...")
        pptx_path = convert_ppt_to_pptx(file_path)
        if pptx_path:
            return extract_text_from_pptx(pptx_path, trigger_ocr, ocr_mix)
        else:
            logging.error("Failed to process .ppt file.")
            return None
    # elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"] and (trigger_ocr or ocr_mix):
    elif ext in Image.registered_extensions() and (trigger_ocr or ocr_mix):
        return extract_text_from_image(file_path)
    else:
        logging.error(f"Unsupported file type: {ext}")
        return None

def save_to_txt(output_path, text):
    """Save extracted text to a .txt file."""
    try:
        with open(output_path, "w", encoding="utf-8") as file:
            file.write(text)
        logging.info(f"Text successfully saved to {output_path}")
    except Exception as e:
        logging.error(f"Failed to save text to {output_path}: {e}")

def process_directory(directory_path, save_all=False, output_file=None, trigger_ocr=False, ocr_mix=False):
    """Process all supported files in the directory recursively."""
    all_text = []
    for root, _, files in os.walk(directory_path):
        for file in files:
            if file.lower().endswith(".txt") and "extracted_texts" in root:
                continue

            file_path = os.path.join(root, file)
            logging.info(f"Processing file: {file_path}")
            text = extract_text_from_file(file_path, trigger_ocr, ocr_mix)
            if text:
                if save_all:
                    all_text.append(f"### {file_path} ###\n{text}\n\n")
                else:
                    output_path = os.path.join(os.path.dirname(file_path), "extracted_texts", os.path.splitext(os.path.basename(file_path))[0] + ".txt")
                    os.makedirs(os.path.dirname(output_path), exist_ok=True)
                    save_to_txt(output_path, text)
    if save_all and output_file:
        save_to_txt(output_file, "\n".join(all_text))

def main():
    if len(sys.argv) < 2 or sys.argv[1] in ("-h", "--help"):
        print("Usage: python main.py <file_or_directory_path> [options]")
        print("       python main.py clear-log")
        print("Options:")
        print("  -h, --help            Show this help message and exit.")
        print("  -a                    Save all extracted text from files in the directory to a single text file (all_extracted_text.txt).")
        print("  --convert pdf         Convert PowerPoint (.ppt or .pptx) to PDF.")
        print("  -ocr, --ocr          Trigger OCR extraction for image files.")
        print("  --ocr-mix             Attempt to extract both regular text and text from images (OCR) in PDF and PPTX files.")
        sys.exit(1)

    if sys.argv[1] == "clear-log":  # Check for 'clear-log' command
            clear_log()
            sys.exit(0)

    path = sys.argv[1]

    if path.startswith("http"):
        scrape_and_save(path)
        sys.exit(0)
    save_all = "-a" in sys.argv
    convert_pdf = "--convert" in sys.argv and "pdf" in sys.argv
    trigger_ocr = True if "-ocr" in sys.argv or "--ocr" in sys.argv else False
    ocr_mix = "--ocr-mix" in sys.argv

    if not os.path.exists(path):
        logging.error(f"Path not found: {path}")
        sys.exit(1)

    output_file = os.path.join(os.path.dirname(path), "extracted_texts", "all_extracted_text.txt") if save_all else None

    if convert_pdf:
        # Convert to PDF if --convert pdf is passed
        if path.lower().endswith(('.ppt', '.pptx')):
            logging.info(f"Converting PowerPoint file to PDF: {path}")
            if path.lower().endswith('.ppt'):
                # Convert .ppt to .pptx first if it's a .ppt file
                pptx_path = convert_ppt_to_pptx(path)
                if pptx_path:
                    path = pptx_path
            pdf_path = convert_pptx_to_pdf(path)
            if pdf_path:
                logging.info(f"Converted to PDF: {pdf_path}")
                sys.exit(0)
            else:
                logging.error(f"Conversion failed for: {path}")
                sys.exit(1)
        else:
            logging.error("Invalid file type. Only .ppt or .pptx files can be converted to PDF.")
            sys.exit(1)

    if os.path.isdir(path):
        logging.info(f"Processing directory: {path}")
        output_folder = os.path.join(path, "extracted_texts")
        os.makedirs(output_folder, exist_ok=True)
        process_directory(path, save_all, output_file, trigger_ocr, ocr_mix)
    else:
        logging.info(f"Processing single file: {path}")
        text = extract_text_from_file(path, trigger_ocr, ocr_mix)
        if text:
            if save_all:
                os.makedirs(os.path.dirname(output_file), exist_ok=True)
                save_to_txt(output_file, text)
            else:
                output_path = os.path.join(os.path.dirname(path), "extracted_texts", os.path.splitext(os.path.basename(path))[0] + ".txt")
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                save_to_txt(output_path, text)


if __name__ == "__main__":
    print("\nTextNomNom - 0xQan\n")
    print("V1: ARCHIVED!!! Update NOW!!! https://github.com/FurqanHun/textnomnom-py")
    main()
    print("\nExitting...\n")
