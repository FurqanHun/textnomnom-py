# app/file_handlers/pptx_handler.py

import io
import logging
from pptx import Presentation
from PIL import Image
import pytesseract


def extract_text_from_pptx(file_path, trigger_ocr=False, ocr_mix=False, callback=None):
    """Extracts text from a .pptx file, with an option for OCR."""
    text_content = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                elif (trigger_ocr or ocr_mix) and hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = pytesseract.image_to_string(
                            Image.open(io.BytesIO(image_bytes))
                        )
                        if ocr_text:
                            slide_text.append(
                                f"[OCR from Slide {i}]\n{ocr_text.strip()}"
                            )
                    except Exception as e:
                        logging.warning(f"OCR on slide {i} failed: {e}")
            if slide_text:
                text_content.append(f"[Slide {i}]\n" + "\n".join(slide_text))
            if callback:
                callback()
        logging.info(f"Extracted text from PowerPoint: {file_path}")
        return "\n\n".join(text_content)
    except Exception as e:
        logging.error(f"Could not read PowerPoint file {file_path}: {e}")
        return None
